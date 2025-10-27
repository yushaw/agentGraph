"""Document content extractors for PDF, DOCX, XLSX, PPTX.

提供各种文档类型的文本提取功能，支持预览和完整提取。

分块策略优化（2025-10-27）:
- Chunk size: 400 字符（约 100-130 tokens 中文，业界最佳实践）
- Chunk overlap: 80 字符（20% 重叠）
- 内容感知分块：尊重段落边界，避免在句子中间截断
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Dict, List

from generalAgent.config.settings import get_settings

LOGGER = logging.getLogger(__name__)

# 支持的文档类型
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}


def _split_text_with_overlap(text: str, max_size: int, overlap: int) -> List[str]:
    """将文本分割为带重叠的块（内容感知）

    策略：
    1. 优先按段落分割（双换行符）
    2. 如果段落太大，按句子分割（句号、问号、感叹号）
    3. 如果句子还是太大，按固定长度分割并添加重叠

    Args:
        text: 原始文本
        max_size: 最大块大小（字符数）
        overlap: 重叠大小（字符数）

    Returns:
        分割后的文本块列表
    """
    if len(text) <= max_size:
        return [text]

    chunks = []

    # 首先按段落分割（双换行或单换行）
    paragraphs = re.split(r'\n\s*\n|\n', text)

    current_chunk = []
    current_length = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        para_len = len(para)

        # 如果单个段落超过 max_size，需要进一步拆分
        if para_len > max_size:
            # 保存当前块
            if current_chunk:
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = []
                current_length = 0

            # 拆分大段落
            sub_chunks = _split_large_paragraph(para, max_size, overlap)
            chunks.extend(sub_chunks)

        # 如果加上当前段落会超过限制
        elif current_length + para_len > max_size:
            # 保存当前块
            if current_chunk:
                chunks.append('\n\n'.join(current_chunk))

            # 计算重叠部分
            overlap_text = _get_overlap_text(current_chunk, overlap)
            current_chunk = [overlap_text, para] if overlap_text else [para]
            current_length = len('\n\n'.join(current_chunk))

        else:
            # 加入当前块
            current_chunk.append(para)
            current_length += para_len + 2  # +2 for '\n\n'

    # 保存最后一个块
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))

    return [c.strip() for c in chunks if c.strip()]


def _split_large_paragraph(paragraph: str, max_size: int, overlap: int) -> List[str]:
    """拆分大段落（按句子或固定长度）"""
    # 尝试按句子分割（中文和英文）
    sentences = re.split(r'([。！？\.!?])', paragraph)

    # 重新组合句子和标点
    reformed_sentences = []
    for i in range(0, len(sentences) - 1, 2):
        if i + 1 < len(sentences):
            reformed_sentences.append(sentences[i] + sentences[i + 1])
        else:
            reformed_sentences.append(sentences[i])

    # 如果没有成功分割，回退到固定长度
    if len(reformed_sentences) <= 1:
        return _split_fixed_size(paragraph, max_size, overlap)

    # 合并短句子
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in reformed_sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        sent_len = len(sentence)

        if sent_len > max_size:
            # 单个句子就超过限制，固定长度切分
            if current_chunk:
                chunks.append(''.join(current_chunk))
                current_chunk = []
                current_length = 0

            chunks.extend(_split_fixed_size(sentence, max_size, overlap))

        elif current_length + sent_len > max_size:
            if current_chunk:
                chunks.append(''.join(current_chunk))

            # 添加重叠
            overlap_text = _get_overlap_text(current_chunk, overlap)
            current_chunk = [overlap_text, sentence] if overlap_text else [sentence]
            current_length = len(''.join(current_chunk))

        else:
            current_chunk.append(sentence)
            current_length += sent_len

    if current_chunk:
        chunks.append(''.join(current_chunk))

    return [c.strip() for c in chunks if c.strip()]


def _split_fixed_size(text: str, max_size: int, overlap: int) -> List[str]:
    """按固定大小分割文本（带重叠）"""
    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + max_size, text_len)
        chunk = text[start:end].strip()

        if chunk:
            chunks.append(chunk)

        # 下一个块的起始位置（考虑重叠）
        start = end - overlap if end < text_len else text_len

    return chunks


def _get_overlap_text(chunks: List[str], overlap_size: int) -> str:
    """从之前的块中提取重叠文本"""
    if not chunks:
        return ""

    combined = '\n\n'.join(chunks)
    if len(combined) <= overlap_size:
        return combined

    return combined[-overlap_size:]


def get_document_info(file_path: Path) -> Dict:
    """获取文档基本信息（页数、字符数等）"""
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _get_pdf_info(file_path)
    elif ext == ".docx":
        return _get_docx_info(file_path)
    elif ext == ".xlsx":
        return _get_xlsx_info(file_path)
    elif ext == ".pptx":
        return _get_pptx_info(file_path)
    else:
        return {"pages": 0, "chars": 0, "type": "unknown"}


def extract_preview(file_path: Path, max_pages: int, max_chars: int) -> str:
    """提取文档预览（前N页/sheets/slides）"""
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf_preview(file_path, max_pages, max_chars)
        elif ext == ".docx":
            return _extract_docx_preview(file_path, max_pages, max_chars)
        elif ext == ".xlsx":
            return _extract_xlsx_preview(file_path, max_pages, max_chars)  # max_pages = sheets
        elif ext == ".pptx":
            return _extract_pptx_preview(file_path, max_pages, max_chars)  # max_pages = slides
        else:
            return f"Error: Unsupported document type: {ext}"
    except Exception as e:
        LOGGER.error(f"Failed to extract preview from {file_path}: {e}")
        return f"Error: Failed to extract preview - {str(e)}"


def extract_full_document(file_path: Path) -> str:
    """提取文档完整内容"""
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf_full(file_path)
        elif ext == ".docx":
            return _extract_docx_full(file_path)
        elif ext == ".xlsx":
            return _extract_xlsx_full(file_path)
        elif ext == ".pptx":
            return _extract_pptx_full(file_path)
        else:
            return f"Error: Unsupported document type: {ext}"
    except Exception as e:
        LOGGER.error(f"Failed to extract content from {file_path}: {e}")
        return f"Error: Failed to extract content - {str(e)}"


def chunk_document(file_path: Path) -> List[Dict]:
    """将文档分块（用于索引）

    Returns:
        List of chunks, each with:
        - id: chunk 编号
        - page: 页码/sheet编号/slide编号
        - text: 文本内容
        - offset: 字符偏移量
    """
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _chunk_pdf(file_path)
        elif ext == ".docx":
            return _chunk_docx(file_path)
        elif ext == ".xlsx":
            return _chunk_xlsx(file_path)
        elif ext == ".pptx":
            return _chunk_pptx(file_path)
        else:
            return []
    except Exception as e:
        LOGGER.error(f"Failed to chunk document {file_path}: {e}")
        return []


# ============ PDF 处理 ============

def _get_pdf_info(file_path: Path) -> Dict:
    """获取 PDF 信息"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    total_chars = sum(len(page.get_text()) for page in doc)
    result = {
        "pages": len(doc),
        "chars": total_chars,
        "type": "pdf"
    }
    doc.close()
    return result


def _extract_pdf_preview(file_path: Path, max_pages: int, max_chars: int) -> str:
    """提取 PDF 前 N 页"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages_content = []
    total_chars = 0

    for i in range(min(max_pages, len(doc))):
        page = doc[i]
        text = page.get_text()

        # 检查是否超过字符限制
        if total_chars + len(text) > max_chars:
            # 截断当前页
            remaining = max_chars - total_chars
            if remaining > 0:
                pages_content.append(f"=== Page {i+1} ===\n{text[:remaining]}...")
            break

        pages_content.append(f"=== Page {i+1} ===\n{text}")
        total_chars += len(text)

    doc.close()
    return "\n\n".join(pages_content)


def _extract_pdf_full(file_path: Path) -> str:
    """提取 PDF 完整内容"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages_content = []

    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            pages_content.append(f"=== Page {i} ===\n{text}")

    doc.close()
    return "\n\n".join(pages_content)


def _chunk_pdf(file_path: Path) -> List[Dict]:
    """PDF 分块（内容感知 + 重叠）

    策略：
    1. 提取每一页的文本
    2. 使用内容感知分块（按段落/句子）
    3. 添加 20% 重叠
    4. 保留页码信息用于引用
    """
    import fitz  # PyMuPDF

    settings = get_settings()
    max_size = settings.documents.chunk_max_size
    overlap = settings.documents.chunk_overlap
    min_size = settings.documents.chunk_min_size

    chunks = []
    chunk_id = 0

    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, 1):
        # 提取文本
        text = page.get_text()

        if not text.strip():
            continue

        # 使用内容感知分块
        page_chunks = _split_text_with_overlap(text, max_size, overlap)

        for chunk_text in page_chunks:
            if len(chunk_text) >= min_size:
                chunks.append({
                    "id": chunk_id,
                    "page": page_num,
                    "text": chunk_text,
                    "offset": sum(len(c["text"]) for c in chunks)
                })
                chunk_id += 1

    doc.close()
    return chunks


# ============ DOCX 处理 ============

def _get_docx_info(file_path: Path) -> Dict:
    """获取 DOCX 信息"""
    from docx import Document

    doc = Document(file_path)
    total_chars = sum(len(para.text) for para in doc.paragraphs)

    return {
        "pages": max(1, len(doc.paragraphs) // 10),  # 估算页数
        "chars": total_chars,
        "type": "docx"
    }


def _extract_docx_preview(file_path: Path, max_pages: int, max_chars: int) -> str:
    """提取 DOCX 前 N 段（估算为页）"""
    from docx import Document

    doc = Document(file_path)
    paragraphs = []
    total_chars = 0

    # 每页约 10 段
    max_paragraphs = max_pages * 10

    for para in doc.paragraphs[:max_paragraphs]:
        text = para.text.strip()
        if not text:
            continue

        if total_chars + len(text) > max_chars:
            remaining = max_chars - total_chars
            if remaining > 0:
                paragraphs.append(text[:remaining] + "...")
            break

        paragraphs.append(text)
        total_chars += len(text)

    return "\n\n".join(paragraphs)


def _extract_docx_full(file_path: Path) -> str:
    """提取 DOCX 完整内容"""
    from docx import Document

    doc = Document(file_path)
    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]

    return "\n\n".join(paragraphs)


def _chunk_docx(file_path: Path) -> List[Dict]:
    """DOCX 分块（内容感知 + 重叠）

    策略：
    1. 提取所有段落
    2. 使用内容感知分块（尊重段落边界）
    3. 添加 20% 重叠
    4. 估算页码用于引用
    """
    from docx import Document

    settings = get_settings()
    max_size = settings.documents.chunk_max_size
    overlap = settings.documents.chunk_overlap
    min_size = settings.documents.chunk_min_size

    doc = Document(file_path)

    # 收集所有文本
    full_text = "\n\n".join(para.text.strip() for para in doc.paragraphs if para.text.strip())

    if not full_text:
        return []

    # 使用内容感知分块
    text_chunks = _split_text_with_overlap(full_text, max_size, overlap)

    chunks = []
    for i, chunk_text in enumerate(text_chunks):
        if len(chunk_text) >= min_size:
            # 估算页码（假设每 2000 字符一页）
            page_estimate = i // 5 + 1

            chunks.append({
                "id": i,
                "page": page_estimate,
                "text": chunk_text,
                "offset": sum(len(c["text"]) for c in chunks)
            })

    return chunks


# ============ XLSX 处理 ============

def _get_xlsx_info(file_path: Path) -> Dict:
    """获取 XLSX 信息"""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    sheet_count = len(wb.sheetnames)

    # 估算字符数
    total_chars = 0
    for sheet_name in wb.sheetnames[:3]:  # 只统计前3个sheet
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True, max_row=100):
            total_chars += sum(len(str(cell)) for cell in row if cell is not None)

    wb.close()

    return {
        "pages": sheet_count,  # Sheet 数量作为页数
        "chars": total_chars,
        "type": "xlsx"
    }


def _extract_xlsx_preview(file_path: Path, max_sheets: int, max_chars: int) -> str:
    """提取 XLSX 前 N 个 sheet"""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    sheets_content = []
    total_chars = 0

    for sheet_name in wb.sheetnames[:max_sheets]:
        sheet = wb[sheet_name]

        rows = []
        for row in sheet.iter_rows(values_only=True, max_row=100):  # 最多100行
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)
                total_chars += len(row_text)

                if total_chars > max_chars:
                    break

        if rows:
            sheet_content = f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows)
            sheets_content.append(sheet_content)

        if total_chars > max_chars:
            break

    wb.close()
    return "\n\n".join(sheets_content)


def _extract_xlsx_full(file_path: Path) -> str:
    """提取 XLSX 完整内容"""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    sheets_content = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        rows = []
        for row in sheet.iter_rows(values_only=True, max_row=1000):  # 最多1000行
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)

        if rows:
            sheet_content = f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows)
            sheets_content.append(sheet_content)

    wb.close()
    return "\n\n".join(sheets_content)


def _chunk_xlsx(file_path: Path) -> List[Dict]:
    """XLSX 分块（按行批次 + 重叠）

    策略：
    1. 遍历所有 sheets
    2. 每个 sheet 按 N 行一批进行分块（默认 20 行）
    3. 添加行重叠（overlap 换算为行数）
    4. 保留 sheet 名称和行号用于引用
    """
    import openpyxl

    settings = get_settings()
    rows_per_chunk = settings.documents.xlsx_rows_per_chunk  # 默认 20 行
    chunk_overlap_rows = max(1, settings.documents.chunk_overlap // 50)  # 约 2-3 行重叠

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    chunks = []
    chunk_id = 0

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # 收集所有行
        all_rows = []
        for row in sheet.iter_rows(values_only=True, max_row=1000):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                all_rows.append(row_text)

        if not all_rows:
            continue

        # 按批次分块（带重叠）
        start_row = 0
        while start_row < len(all_rows):
            end_row = min(start_row + rows_per_chunk, len(all_rows))
            chunk_rows = all_rows[start_row:end_row]

            if chunk_rows:
                row_range = f"rows {start_row + 1}-{end_row}"
                chunk_text = f"[Sheet: {sheet_name}, {row_range}]\n" + "\n".join(chunk_rows)

                chunks.append({
                    "id": chunk_id,
                    "page": chunk_id + 1,  # 使用 chunk_id 作为"页码"
                    "text": chunk_text,
                    "offset": sum(len(c["text"]) for c in chunks)
                })
                chunk_id += 1

            # 下一批的起始位置（考虑重叠）
            start_row = end_row - chunk_overlap_rows if end_row < len(all_rows) else len(all_rows)

    wb.close()
    return chunks


# ============ PPTX 处理 ============

def _get_pptx_info(file_path: Path) -> Dict:
    """获取 PPTX 信息"""
    from pptx import Presentation

    prs = Presentation(file_path)
    total_chars = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                total_chars += len(shape.text)

    return {
        "pages": len(prs.slides),  # Slide 数量作为页数
        "chars": total_chars,
        "type": "pptx"
    }


def _extract_pptx_preview(file_path: Path, max_slides: int, max_chars: int) -> str:
    """提取 PPTX 前 N 张幻灯片"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_content = []
    total_chars = 0

    for slide_idx, slide in enumerate(prs.slides[:max_slides], 1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
                        total_chars += len(text)
            except Exception as e:
                # Skip shapes that can't be read
                LOGGER.debug(f"Failed to read shape in slide {slide_idx}: {e}")
                continue

        if texts:
            slide_content = f"=== Slide {slide_idx} ===\n" + "\n\n".join(texts)
            slides_content.append(slide_content)

        if total_chars > max_chars:
            break

    return "\n\n".join(slides_content)


def _extract_pptx_full(file_path: Path) -> str:
    """提取 PPTX 完整内容"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_content = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            except Exception as e:
                LOGGER.debug(f"Failed to read shape in slide {slide_idx}: {e}")
                continue

        if texts:
            slide_content = f"=== Slide {slide_idx} ===\n" + "\n\n".join(texts)
            slides_content.append(slide_content)

    return "\n\n".join(slides_content)


def _chunk_pptx(file_path: Path) -> List[Dict]:
    """PPTX 分块（按 slide + 大 slide 拆分）

    策略：
    1. 通常按 slide 分块（一个 slide 一个 chunk）
    2. 如果 slide 文本超过 max_size，使用内容感知分块
    3. 保留 slide 编号用于引用
    """
    from pptx import Presentation

    settings = get_settings()
    max_size = settings.documents.chunk_max_size
    overlap = settings.documents.chunk_overlap
    min_size = settings.documents.chunk_min_size

    prs = Presentation(file_path)
    chunks = []
    chunk_id = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            except Exception as e:
                LOGGER.debug(f"Failed to read shape in slide {slide_idx}: {e}")
                continue

        if not texts:
            continue

        slide_text = "\n\n".join(texts)

        # 如果 slide 内容不大，整个 slide 作为一个 chunk
        if len(slide_text) <= max_size:
            chunks.append({
                "id": chunk_id,
                "page": slide_idx,
                "text": slide_text,
                "offset": sum(len(c["text"]) for c in chunks)
            })
            chunk_id += 1
        else:
            # 大 slide 需要拆分
            slide_chunks = _split_text_with_overlap(slide_text, max_size, overlap)

            for sub_chunk in slide_chunks:
                if len(sub_chunk) >= min_size:
                    chunks.append({
                        "id": chunk_id,
                        "page": slide_idx,
                        "text": sub_chunk,
                        "offset": sum(len(c["text"]) for c in chunks)
                    })
                    chunk_id += 1

    return chunks
